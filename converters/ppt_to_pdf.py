from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
import os
import time
import textwrap

def convert_ppt_to_pdf(input_path, output_folder):
    try:
        prs = Presentation(input_path)

        os.makedirs(output_folder, exist_ok=True)
        output_path = os.path.join(
            output_folder, f"converted_{int(time.time())}.pdf"
        )

        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter

        for slide_num, slide in enumerate(prs.slides, start=1):

            # ================= BACKGROUND =================
            c.setFillColorRGB(0.92, 0.92, 1)
            c.rect(0, 0, width, height, fill=1)

            c.setFillColorRGB(0.85, 0.85, 1)
            c.rect(0, height/2, width, height/2, fill=1)

            # ================= CARD =================
            card_x = 50
            card_y = 80
            card_w = width - 100
            card_h = height - 160

            # shadow
            c.setFillColorRGB(0.8, 0.8, 0.9)
            c.roundRect(card_x + 6, card_y - 6, card_w, card_h, 18, fill=1)

            # white card
            c.setFillColorRGB(1, 1, 1)
            c.roundRect(card_x, card_y, card_w, card_h, 18, fill=1)

            # border
            c.setStrokeColorRGB(0.55, 0.25, 0.85)
            c.setLineWidth(3)
            c.roundRect(card_x, card_y, card_w, card_h, 18, fill=0)

            # header
            c.setFillColorRGB(0.55, 0.25, 0.85)
            c.roundRect(card_x, card_y + card_h - 55, card_w, 55, 18, fill=1)

            # ================= EXTRACT TEXT =================
            title_text = None
            content_lines = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            if not title_text:
                                title_text = text
                            else:
                                content_lines.append(text)

            # ================= TITLE =================
            c.setFillColorRGB(1, 1, 1)
            c.setFont("Helvetica-Bold", 18)

            title_y = card_y + card_h - 35

            for line in textwrap.wrap(title_text or "", 40):
                c.drawString(card_x + 20, title_y, line)
                title_y -= 22

            # ================= BODY =================
            y = title_y - 20

            c.setFillColor(colors.black)
            c.setFont("Helvetica", 13)

            for text in content_lines:
                lines = textwrap.wrap(text, 65)

                for line in lines:
                    if y < card_y + 60:
                        break

                    c.drawString(card_x + 35, y, "• " + line)
                    y -= 18

                y -= 6

            # ================= FOOTER =================
            c.setStrokeColorRGB(0.85, 0.85, 0.9)
            c.line(card_x, card_y + 30, card_x + card_w, card_y + 30)

            c.setFont("Helvetica-Bold", 10)
            c.setFillColor(colors.grey)
            c.drawRightString(width - 40, 30, f"Slide {slide_num}")

            c.showPage()

        c.save()
        return output_path

    except Exception as e:
        raise Exception(f"PPT to PDF conversion failed: {str(e)}")