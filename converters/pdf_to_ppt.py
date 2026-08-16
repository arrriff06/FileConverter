import os
import fitz  # PyMuPDF
from pptx import Presentation
import time


def convert_pdf_to_ppt(pdf_path, output_folder):
    try:
        pdf_path = os.path.abspath(pdf_path)
        output_folder = os.path.abspath(output_folder)

        os.makedirs(output_folder, exist_ok=True)

        # Open PDF
        doc = fitz.open(pdf_path)

        # Create PPT
        prs = Presentation()

        temp_images = []

        for i, page in enumerate(doc):
            # Convert page to image
            pix = page.get_pixmap()
            img_path = os.path.join(output_folder, f"temp_{i}_{int(time.time())}.png")
            pix.save(img_path)
            temp_images.append(img_path)

            # Add slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                img_path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

        # Save output
        output_file = os.path.join(output_folder, f"converted_{int(time.time())}.pptx")
        prs.save(output_file)

        doc.close()

        # Cleanup temp images
        for img in temp_images:
            if os.path.exists(img):
                os.remove(img)

        return output_file

    except Exception as e:
        raise Exception(f"PDF to PPT conversion failed: {str(e)}")